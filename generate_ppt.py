"""Generate ANISAS project PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG_DARK  = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD  = RGBColor(0x1E, 0x29, 0x3B)
ACCENT   = RGBColor(0x3B, 0x82, 0xF6)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
RED      = RGBColor(0xEF, 0x44, 0x44)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
YELLOW   = RGBColor(0xEA, 0xB3, 0x08)
PURPLE   = RGBColor(0xA8, 0x55, 0xF7)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x94, 0xA3, 0xB8)
BORDER   = RGBColor(0x33, 0x41, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide, color=BG_DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def shape(slide, l, t, w, h, fill=None, line=None, lw=1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = b
    p.font.name = "Calibri"
    p.alignment = a
    return tb


def bullets(slide, l, t, w, h, items, sz=15, c=WHITE, sp=8):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.name = "Calibri"
        p.space_after = Pt(sp)
    return tb


def card(slide, l, t, w, h, title, body, tc=ACCENT, bc=MUTED):
    shape(slide, l, t, w, h, fill=BG_CARD, line=BORDER)
    txt(slide, l+Inches(0.2), t+Inches(0.15), w-Inches(0.4), Inches(0.4), title, 14, tc, True)
    txt(slide, l+Inches(0.2), t+Inches(0.5), w-Inches(0.4), h-Inches(0.6), body, 12, bc)


def title_bar(slide, text, color=ACCENT):
    txt(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.7), text, 34, color, True)
    shape(slide, Inches(0.8), Inches(1.05), Inches(4), Pt(2), fill=color)


def cmd_line(slide, text, y=5.8):
    txt(slide, Inches(0.8), Inches(y), Inches(11.7), Inches(0.5), text, 15, GREEN, True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 - Title
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
shape(s, Inches(4), Inches(2.0), Inches(5.3), Pt(3), fill=ACCENT)
txt(s, Inches(1), Inches(2.3), Inches(11.3), Inches(1.2), "ANISAS", 54, ACCENT, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8), "Autonomous Network Intelligence\n& Security Assessment System", 28, WHITE, False, PP_ALIGN.CENTER)
shape(s, Inches(4), Inches(4.4), Inches(5.3), Pt(3), fill=ACCENT)
txt(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5), "A 7-Module Python Platform for Automated Network Security Intelligence", 16, MUTED, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.5), "Version 1.0.0  |  Python 3.10+", 14, MUTED, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 - What is ANISAS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "What is ANISAS?")
txt(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1), "ANISAS is a fully automated network security intelligence platform that takes\na single IP address and produces a complete security assessment.", 18, WHITE)
bullets(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(4.5), [
    "Accepts a target IP address and runs 6 analysis modules automatically",
    "Covers ASN lookup, port scanning, firewall detection, IoT fingerprinting, wireless analysis",
    "Uses AI/ML decision trees for device classification and risk scoring",
    "Generates professional PDF reports and interactive web dashboard",
    "Built entirely in Python - no external paid services required",
    "Designed for security researchers, network admins, and SOC teams",
], 16, WHITE, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 - Problem Statement
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Problem Statement")
problems = [
    ("Manual Security Audits", "Traditional network audits require multiple tools and hours of manual analysis"),
    ("Tool Fragmentation", "Security teams juggle 5-10 different tools for a single assessment"),
    ("No Unified View", "Results are scattered across different formats and dashboards"),
    ("Slow Response Time", "By the time reports are compiled, vulnerabilities may already be exploited"),
    ("Skill Gap", "Not every organization has dedicated security analysts"),
]
for i, (t, d) in enumerate(problems):
    y = Inches(1.5) + Inches(1.1) * i
    shape(s, Inches(0.8), y, Inches(11.7), Inches(0.95), fill=BG_CARD, line=BORDER)
    txt(s, Inches(1.1), y + Inches(0.08), Inches(5), Inches(0.4), t, 16, ORANGE, True)
    txt(s, Inches(1.1), y + Inches(0.45), Inches(11), Inches(0.4), d, 13, MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 - Our Solution
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Our Solution", GREEN)
txt(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.7), "One IP in  ->  Complete security intelligence out", 20, WHITE, True)
for i, (t, b) in enumerate([
    ("Single Command", "Run one command\nor click one button\nin the dashboard"),
    ("Automated Pipeline", "6 modules run\nsequentially, each\nfeeding the next"),
    ("AI-Powered Analysis", "Decision trees classify\ndevices and score\nrisks automatically"),
]):
    card(s, Inches(0.8) + Inches(4.1) * i, Inches(2.5), Inches(3.7), Inches(2), t, b, GREEN, WHITE)
txt(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1), "Output: Professional PDF report, interactive web dashboard with live topology graph,\nJSON data for integration with SIEM/SOAR tools", 16, MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 - Architecture
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "System Architecture")
mods = [
    ("Module 1\nASN & ISP", ACCENT, 0.5, 1.8),
    ("Module 2\nNetwork Recon", CYAN, 0.5, 3.0),
    ("Module 3\nPerimeter", RED, 0.5, 4.2),
    ("Module 4\nIoT/Surveillance", ORANGE, 4.5, 1.8),
    ("Module 5\nWireless", PURPLE, 4.5, 3.0),
    ("Module 6\nAI/ML Engine", GREEN, 4.5, 4.2),
]
for label, color, x, y in mods:
    sh = shape(s, Inches(x), Inches(y), Inches(3.5), Inches(0.95), fill=BG_CARD, line=color, lw=2)
    txt(s, Inches(x+0.2), Inches(y+0.12), Inches(3.1), Inches(0.75), label, 14, color, True, PP_ALIGN.CENTER)
shape(s, Inches(8.8), Inches(1.8), Inches(4), Inches(3.35), fill=BG_CARD, line=ACCENT, lw=2)
txt(s, Inches(9.0), Inches(1.95), Inches(3.6), Inches(0.4), "Module 7", 18, ACCENT, True, PP_ALIGN.CENTER)
txt(s, Inches(9.0), Inches(2.4), Inches(3.6), Inches(0.4), "GUI Dashboard", 16, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(9.0), Inches(2.9), Inches(3.6), Inches(2), "FastAPI Backend\nSSE Live Streaming\nvis-network Topology\nPDF/JSON Export\nVanilla JS Frontend", 13, MUTED, False, PP_ALIGN.CENTER)
for y in [2.2, 3.4, 4.6]:
    txt(s, Inches(4.0), Inches(y), Inches(0.5), Inches(0.5), "->", 24, MUTED, True)
txt(s, Inches(8.3), Inches(3.2), Inches(0.5), Inches(0.5), "->", 24, MUTED, True)

# ═══════════════════════════════════════════════════════════════
# SLIDES 6-12 - Module Details
# ═══════════════════════════════════════════════════════════════
# Module 1
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 1 - ASN & ISP Intelligence")
txt(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5), "Maps the complete ASN/ISP ecosystem for any public IP address", 18, WHITE)
bullets(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.5), [
    "ASN Resolution: ipinfo.io -> bgpview.io -> Team Cymru",
    "ISP Profile: Name, NOC contact, abuse contact, peering partners",
    "Multi-ASN Detection: Identifies primary and secondary ASNs",
    "IP Prefix Enumeration: Discovers all CIDR blocks",
], 14, WHITE, 10)
bullets(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.5), [
    "AI/NLP Risk Analysis using Hugging Face transformers",
    "Sentiment-analysis model maps metadata to risk levels",
    "Keyword fallback when model is unavailable",
    "Generates JSON + PDF output reports",
], 14, WHITE, 10)
cmd_line(s, "Command:  python -m anisas 8.8.8.8")

# Module 2
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 2 - Network Reconnaissance", CYAN)
txt(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5), "Discovers hosts, scans ports, fingerprints OS, and builds topology", 18, WHITE)
bullets(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5), [
    "Subnet Enumeration: Discovers CIDR blocks and VLAN tags",
    "Host Discovery: ARP, ICMP, and TCP-SYN probes with configurable threads",
    "Port Scanning: Top 1000 TCP/UDP ports with service and banner detection",
    "OS Fingerprinting: TTL analysis + TCP window size matching",
    "Topology Graph: Builds node-edge graph for dashboard visualization",
    "Stealth Mode: Randomized timing, fragmented packets, rate limiting",
], 14, WHITE, 10)
cmd_line(s, "Command:  python -m anisas.recon 192.168.1.0/24")

# Module 3
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 3 - Security Perimeter Detection", RED)
for i, (t, b, c) in enumerate([
    ("Firewall", "Type: Stateful/Stateless\nBehavior: Filtered/Unfiltered", RED),
    ("IDS/IPS", "Action: TCP-RST, Drop\nICMP-Unreachable", ORANGE),
    ("DMZ", "Boundary: Public-DMZ\nInternal-Only, Hybrid", YELLOW),
    ("WAF", "Vendor identification\nSignature matching", PURPLE),
]):
    x = Inches(0.8) + Inches(3.1) * i
    sh = shape(s, x, Inches(1.4), Inches(2.8), Inches(1.8), fill=BG_CARD, line=c, lw=2)
    txt(s, x+Inches(0.15), Inches(1.5), Inches(2.5), Inches(0.35), t, 15, c, True, PP_ALIGN.CENTER)
    txt(s, x+Inches(0.15), Inches(1.9), Inches(2.5), Inches(1.1), b, 12, MUTED, False, PP_ALIGN.CENTER)
bullets(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(2), [
    "Evasion Benchmarks: Fragmentation testing, slow-rate timing analysis",
    "AI Detection Prediction: ML model predicts probe detection probability",
    "Overall Posture: Aggregated risk level (Low/Medium/High) with summary",
], 14, WHITE, 10)
cmd_line(s, "Command:  python -m anisas.perimeter 192.168.1.1")

# Module 4
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 4 - Surveillance & IoT Fingerprinting", ORANGE)
txt(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5), "Identifies cameras, NVRs, DVRs, and IoT devices on the network", 18, WHITE)
bullets(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5), [
    "Protocol Scanning: RTSP (554), HTTP (80/443), ONVIF (80/8080)",
    "OUI Lookup: MAC prefix -> vendor identification (Hikvision, Dahua, etc.)",
    "Device Classification: CCTV Camera, NVR, DVR, XVR, Generic IoT",
    "Firmware Detection: Extracts version strings from HTTP banners",
    "CVE Cross-Reference: Maps vendor/model to known vulnerabilities",
    "IP Range Prediction: ML-based prediction of additional IoT IP ranges",
    "Risk Rating: CRITICAL / HIGH / MEDIUM / LOW per device",
], 14, WHITE, 8)
cmd_line(s, "Command:  python -m anisas.iot 192.168.1.0/24")

# Module 5
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 5 - Wireless Network Intelligence", PURPLE)
bullets(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.5), [
    "AP Enumeration: SSID, BSSID, channel, RSSI, encryption type, vendor OUI",
    "Client Discovery: MAC address, assigned IP, hostname, active/inactive status",
    "Authentication Analysis: Primary method, MAC filtering detection",
    "MAC Cloning PoC: Demonstrates cloning inactive MAC to gain access (dry run)",
    "AI Anomaly Detection: Clusters devices and flags anomalous MACs",
    "Hardening Recommendations: Actionable security guidance",
], 14, WHITE, 12)
txt(s, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5), "Note: Requires wireless hardware and admin privileges", 14, ORANGE, True)
cmd_line(s, "Command:  python -m anisas.wireless")

# Module 6
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 6 - AI/ML Analytics Engine", GREEN)
txt(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5), "Aggregates all module data and applies ML classification + risk scoring", 18, WHITE)
headers = ["Capability", "Method", "Accuracy"]
rows = [
    ["Device Classification", "Decision tree on port/OS/TTL features", "~95%"],
    ["OS Fingerprinting", "Decision tree on TTL/window size/ports", "~83%"],
    ["Topology Inference", "Graph analysis of host-subnet relationships", "Heuristic"],
    ["Anomaly Detection", "Statistical outlier detection on risk scores", "Threshold"],
    ["Risk Scoring", "Weighted formula (device + OS + ports + anomalies)", "0-10 scale"],
    ["NL Summary", "Template-based executive summary generation", "N/A"],
]
y0 = Inches(2.0)
cx = [Inches(0.8), Inches(4.3), Inches(9.8)]
cw = [Inches(3.5), Inches(5.5), Inches(2.5)]
shape(s, Inches(0.8), y0, Inches(11.7), Inches(0.45), fill=ACCENT)
for j, h in enumerate(headers):
    txt(s, cx[j], y0+Inches(0.05), cw[j], Inches(0.35), h, 13, WHITE, True, PP_ALIGN.CENTER)
for i, row in enumerate(rows):
    y = y0 + Inches(0.45) + Inches(0.5) * i
    bgc = BG_CARD if i % 2 == 0 else RGBColor(0x14, 0x1B, 0x2D)
    shape(s, Inches(0.8), y, Inches(11.7), Inches(0.5), fill=bgc, line=RGBColor(0x1E, 0x2D, 0x4A))
    for j, cell in enumerate(row):
        txt(s, cx[j], y+Inches(0.08), cw[j], Inches(0.35), cell, 12, GREEN if j==2 else WHITE, j==2, PP_ALIGN.CENTER)
cmd_line(s, "Command:  python -m anisas.ai_engine mod1.json mod2.json mod3.json mod4.json mod5.json")

# Module 7
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Module 7 - GUI Dashboard")
txt(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5), "FastAPI web application with real-time scan visualization", 18, WHITE)
bullets(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4), [
    "Live Scan Progress: Real-time module status with animated indicators",
    "SSE Log Streaming: Color-coded log entries with timestamps",
    "Network Topology: Interactive vis-network graph with device coloring",
    "Perimeter Overlay: Firewall/IDS/DMZ/WAF detection status badges",
    "AI Summary Tab: Health score, anomaly count, risk assessment, mitigations",
    "Threats Tab: Key threats with severity badges, device classification list",
    "Node Detail Modal: Click any node for full device details",
    "PDF Export: Multi-page professional report with cover page and tables",
    "JSON Export: Complete machine-readable results for SIEM integration",
], 13, WHITE, 8)
cmd_line(s, "Command:  python -m anisas.dashboard   ->   http://127.0.0.1:8000")

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 - Tech Stack
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Technology Stack")
tech = [
    ("Language", "Python 3.10+", ACCENT),
    ("Data Models", "Pydantic v2 (validation + JSON schema)", CYAN),
    ("HTTP Client", "httpx (async, auto-failover)", GREEN),
    ("PDF Reports", "ReportLab (tables, styles, multi-page)", ORANGE),
    ("NLP/ML", "Hugging Face Transformers + PyTorch", PURPLE),
    ("Web Framework", "FastAPI + Uvicorn (async ASGI)", RED),
    ("Frontend", "Vanilla JS + vis-network (topology graph)", YELLOW),
    ("Streaming", "Server-Sent Events (SSE)", CYAN),
]
for i, (label, desc, color) in enumerate(tech):
    r, c = divmod(i, 2)
    x = Inches(0.8) + Inches(6.2) * c
    y = Inches(1.4) + Inches(1.2) * r
    sh = shape(s, x, y, Inches(5.8), Inches(0.95), fill=BG_CARD, line=color, lw=2)
    txt(s, x+Inches(0.2), y+Inches(0.08), Inches(2.5), Inches(0.35), label, 14, color, True)
    txt(s, x+Inches(0.2), y+Inches(0.45), Inches(5.4), Inches(0.4), desc, 12, MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 - Project Structure
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Project Structure")
bullets(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(4.5), [
    "anisas/                    <- Root package (Module 1)",
    "  recon/                   <- Module 2 - Network Reconnaissance (9 files)",
    "  perimeter/               <- Module 3 - Security Perimeter (9 files)",
    "  iot/                     <- Module 4 - IoT/Surveillance (9 files)",
    "  wireless/                <- Module 5 - Wireless Intelligence (9 files)",
    "  ai_engine/               <- Module 6 - AI/ML Analytics (10 files)",
    "  dashboard/               <- Module 7 - GUI Dashboard (6 files + HTML)",
    "",
    "requirements.txt           <- Python dependencies",
    "pyproject.toml             <- Package metadata + build config",
    "README.md                  <- Full documentation",
], 15, WHITE, 6)
txt(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5), "Total: 54 Python files across 7 modules", 16, ACCENT, True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 - How to Use
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "How to Use")
steps = [
    ("Step 1", "Install dependencies", "pip install -r requirements.txt", ACCENT),
    ("Step 2", "Launch dashboard", "python -m anisas.dashboard", GREEN),
    ("Step 3", "Enter target IP", "Type any public IP in the input field", ORANGE),
    ("Step 4", "Click Scan", "Watch live progress via SSE streaming", PURPLE),
    ("Step 5", "Export results", "Download as JSON or professional PDF report", CYAN),
]
for i, (step, desc, cmd, color) in enumerate(steps):
    y = Inches(1.4) + Inches(1.1) * i
    shape(s, Inches(0.8), y, Inches(1.2), Inches(0.9), fill=color)
    txt(s, Inches(0.8), y+Inches(0.15), Inches(1.2), Inches(0.6), step, 16, WHITE, True, PP_ALIGN.CENTER)
    txt(s, Inches(2.2), y+Inches(0.05), Inches(4), Inches(0.4), desc, 16, WHITE, True)
    txt(s, Inches(2.2), y+Inches(0.45), Inches(9), Inches(0.4), cmd, 14, GREEN)

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 - Sample Output
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Sample Output")
card(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.5),
     "PDF Report Contents",
     "Cover page with target IP and timestamp\n"
     "Section 1: ASN details, IP prefixes, ISP profile, AI risk\n"
     "Section 2: Subnets, hosts, open ports, OS fingerprints\n"
     "Section 3: Firewall/IDS/DMZ/WAF detection, evasion tests\n"
     "Section 4: IoT devices, CVE vulnerabilities, risk ratings\n"
     "Section 5: APs, clients, auth analysis, anomaly detection\n"
     "Section 6: Device classifications, risk scores\n"
     "         Executive summary with threats & mitigations\n"
     "Footer with page numbers and generation timestamp", ACCENT, WHITE)
card(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(4.5),
     "JSON Output Structure",
     "module1: ASN details, ISP profile, risk summary\n"
     "module2: Active hosts, open ports, topology graph\n"
     "module3: Perimeter defenses, evasion benchmarks\n"
     "module4: Surveillance devices, CVE entries\n"
     "module5: Access points, clients, anomaly flags\n"
     "module6: Device classifications, risk scores\n"
     "         Executive NL summary with threats & mitigations\n"
     "All validated via Pydantic models", GREEN, WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 - Advantages
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Advantages", GREEN)
advs = [
    ("Fully Automated", "One command runs all 6 analysis modules end-to-end"),
    ("No External Dependencies", "Pure Python - no paid APIs or proprietary tools"),
    ("AI-Powered", "ML classifiers for device type and OS with 83-95% accuracy"),
    ("Real-time Dashboard", "Live scan progress, interactive topology, color-coded logs"),
    ("Professional Reports", "Multi-page PDF with tables, charts, and executive summary"),
    ("Modular Design", "Each module runs independently or as part of the full pipeline"),
    ("Extensible", "Add new modules, classifiers, or output formats easily"),
    ("Lab Ready", "Designed for security research, education, and SOC operations"),
]
for i, (t, d) in enumerate(advs):
    r, c = divmod(i, 2)
    x = Inches(0.8) + Inches(6.2) * c
    y = Inches(1.4) + Inches(1.3) * r
    shape(s, x, y, Inches(5.8), Inches(1.1), fill=BG_CARD, line=GREEN, lw=1)
    txt(s, x+Inches(0.2), y+Inches(0.08), Inches(5.4), Inches(0.35), t, 15, GREEN, True)
    txt(s, x+Inches(0.2), y+Inches(0.45), Inches(5.4), Inches(0.5), d, 12, MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 - Future Scope
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
title_bar(s, "Future Scope", CYAN)
bullets(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5), [
    "Integration with SIEM platforms (Splunk, ELK, Wazuh) via REST API",
    "Real-time monitoring mode with scheduled scans and alerting",
    "Support for IPv6 target addresses and larger network ranges",
    "Additional ML models: deep learning for traffic pattern analysis",
    "Docker containerization for one-click deployment",
    "Multi-user authentication and role-based access control",
    "Cloud deployment (AWS/Azure) for remote network assessment",
    "Plugin system for custom module development",
    "Integration with threat intelligence feeds (VirusTotal, AbuseIPDB)",
    "Mobile-responsive dashboard for field operations",
], 16, WHITE, 10)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 - Thank You
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
shape(s, Inches(4), Inches(2.5), Inches(5.3), Pt(3), fill=ACCENT)
txt(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1), "Thank You", 48, ACCENT, True, PP_ALIGN.CENTER)
shape(s, Inches(4), Inches(3.9), Inches(5.3), Pt(3), fill=ACCENT)
txt(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6), "ANISAS - Autonomous Network Intelligence & Security Assessment System", 18, MUTED, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5), "Version 1.0.0  |  Python 3.10+  |  54 Files  |  7 Modules", 14, MUTED, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.5), "Questions?", 20, WHITE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════
out = r"C:\Users\PREETANSH GOHIL\Desktop\New folder (6)\ANISAS_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
